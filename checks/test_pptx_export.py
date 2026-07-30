from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation

from artifact_fingerprint import presentation_input_fingerprint
from pptx_export import (
    ASPECT_RATIO,
    build_image_only_pptx,
    inspect_pptx_readiness,
)
from visual_provenance import write_visual_provenance


def _prepared_run(tmp_path: Path, slide_ids: tuple[str, ...] = ("slide_001", "slide_002")) -> Path:
    planning = tmp_path / "planning"
    planning.mkdir(parents=True)
    contract = {
        "version": "visual_contract_v1",
        "slides": [
            {"slide_id": slide_id, "title": f"页面 {index}"}
            for index, slide_id in enumerate(slide_ids, start=1)
        ],
    }
    (planning / "visual_contract.json").write_text(
        json.dumps(contract, ensure_ascii=False),
        encoding="utf-8",
    )
    for index, slide_id in enumerate(slide_ids, start=1):
        slide_dir = tmp_path / "slides" / slide_id
        slide_dir.mkdir(parents=True)
        image_path = slide_dir / "visual_draft.png"
        Image.new("RGB", (640, 360), (230 - index * 10, 235, 250)).save(image_path)
        write_visual_provenance(
            tmp_path,
            slide_id,
            image_path=image_path,
            provider="manual_upload",
            source_type="test",
            source_filename=f"{slide_id}.png",
        )
    return tmp_path


def test_pptx_readiness_accepts_ordered_confirmed_images(tmp_path: Path) -> None:
    run_dir = _prepared_run(tmp_path)

    readiness = inspect_pptx_readiness(run_dir)

    assert readiness["ready"] is True
    assert readiness["slide_ids"] == ["slide_001", "slide_002"]
    assert readiness["slide_count"] == 2
    assert readiness["ready_slide_count"] == 2
    assert not readiness["issues"]
    assert abs((readiness["slides"][0]["width"] / readiness["slides"][0]["height"]) - ASPECT_RATIO) < 0.01


def test_pptx_export_creates_verified_16_by_9_deck(tmp_path: Path) -> None:
    run_dir = _prepared_run(tmp_path)
    progress_events: list[tuple[int, str]] = []

    result = build_image_only_pptx(
        run_dir,
        "presentation_test.pptx",
        title="测试项目",
        progress=lambda value, stage: progress_events.append((value, stage)),
    )

    output = Path(result["path"])
    assert output.is_file()
    assert Path(f"{output}.export.json").is_file()
    deck = Presentation(output)
    assert len(deck.slides) == 2
    assert all(len(slide.shapes) == 1 for slide in deck.slides)
    assert abs((deck.slide_width / deck.slide_height) - ASPECT_RATIO) < 0.000001
    assert progress_events[-1] == (100, "completed")
    assert result["metadata"]["content_mode"] == "full_slide_bitmap"
    assert result["metadata"]["slide_ids"] == ["slide_001", "slide_002"]


def test_pptx_readiness_reports_missing_and_wrong_aspect_images(tmp_path: Path) -> None:
    run_dir = _prepared_run(tmp_path)
    (run_dir / "slides" / "slide_001" / "visual_draft.png").unlink()
    wrong = run_dir / "slides" / "slide_002" / "visual_draft.png"
    Image.new("RGB", (500, 500), "white").save(wrong)

    readiness = inspect_pptx_readiness(run_dir)

    assert readiness["ready"] is False
    codes = {issue["code"] for issue in readiness["issues"]}
    assert "missing_image" in codes
    assert "invalid_aspect_ratio" in codes


def test_presentation_fingerprint_changes_when_a_slide_changes(tmp_path: Path) -> None:
    run_dir = _prepared_run(tmp_path, ("slide_001",))
    before = presentation_input_fingerprint(run_dir)
    image_path = run_dir / "slides" / "slide_001" / "visual_draft.png"
    Image.new("RGB", (640, 360), "white").save(image_path)
    after = presentation_input_fingerprint(run_dir)

    assert before["digest"] != after["digest"]
