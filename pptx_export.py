"""Image-only PPTX readiness checks and atomic export generation."""

from __future__ import annotations

from datetime import datetime
import json
import os
from pathlib import Path
from typing import Any, Callable

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from artifact_fingerprint import presentation_input_fingerprint, sha256_file
from pipeline_lifecycle import write_json_atomic
from project_storage import (
    UnsafeProjectPath,
    presentation_file,
    presentation_sidecar,
    presentations_dir,
    safe_identifier,
    slide_dir,
    slide_file,
    visual_contract_path,
)
from visual_provenance import visual_provenance_status


PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PPTX_EXPORT_VERSION = "image_only_pptx_v1"
SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5
ASPECT_RATIO = 16 / 9
ASPECT_TOLERANCE = 0.01


class PptxReadinessError(ValueError):
    def __init__(self, readiness: dict[str, Any]):
        super().__init__("PPTX export inputs are not ready")
        self.readiness = readiness


def _issue(code: str, message: str, slide_id: str = "") -> dict[str, str]:
    value = {"code": code, "message": message}
    if slide_id:
        value["slide_id"] = slide_id
    return value


def _read_slide_notes(run_dir: str | Path, slide_id: str) -> str:
    """读取某页演讲稿作为 PPT 备注文本。

    优先取可读版 narration.txt（clean_tts_text 已去除 TTS 标记）；
    不存在时回退到 tts_text.txt；两者都缺失则返回空串。
    """
    root = Path(run_dir).resolve()
    slide_dir_path = slide_dir(root, slide_id)
    for name in ("narration.txt", "tts_text.txt"):
        candidate = slide_dir_path / name
        if candidate.is_file():
            try:
                text = candidate.read_text(encoding="utf-8-sig").strip()
            except OSError:
                continue
            if text:
                return text
    return ""


def _read_contract(run_dir: str | Path) -> tuple[dict[str, Any] | None, list[dict[str, str]]]:
    path = visual_contract_path(run_dir)
    try:
        payload = json.loads(path.read_text(encoding="utf-8-sig"))
    except FileNotFoundError:
        return None, [_issue("missing_contract", "尚未生成分镜规划，请先完成步骤 2。")]
    except (OSError, json.JSONDecodeError):
        return None, [_issue("invalid_contract", "分镜规划文件无法读取或格式损坏。")]
    if not isinstance(payload, dict) or not isinstance(payload.get("slides"), list):
        return None, [_issue("invalid_contract", "分镜规划缺少有效的 slides 列表。")]
    if not payload["slides"]:
        return payload, [_issue("empty_contract", "分镜中没有可导出的页面。")]
    return payload, []


def inspect_pptx_readiness(run_dir: str | Path) -> dict[str, Any]:
    root = Path(run_dir).resolve()
    contract, issues = _read_contract(root)
    slide_ids: list[str] = []
    slides: list[dict[str, Any]] = []
    seen: set[str] = set()
    if contract:
        for index, slide in enumerate(contract.get("slides") or [], start=1):
            if not isinstance(slide, dict):
                issues.append(_issue("invalid_slide", f"第 {index} 页分镜不是有效对象。"))
                continue
            slide_id = str(slide.get("slide_id") or "").strip()
            if not slide_id:
                issues.append(_issue("missing_slide_id", f"第 {index} 页缺少 slide_id。"))
                continue
            try:
                safe_identifier(slide_id, label="slide_id")
            except UnsafeProjectPath:
                issues.append(_issue("unsafe_slide_id", "页面 ID 不符合安全规则。", slide_id))
                continue
            if slide_id in seen:
                issues.append(_issue("duplicate_slide_id", "页面 ID 重复。", slide_id))
                continue
            seen.add(slide_id)
            slide_ids.append(slide_id)
            image_path = slide_file(root, slide_id, "visual_draft.png")
            if not image_path.is_file():
                issues.append(_issue("missing_image", "尚未确认当前页面图片。", slide_id))
                continue
            try:
                with Image.open(image_path) as image:
                    image.verify()
                with Image.open(image_path) as image:
                    width, height = image.size
            except (OSError, ValueError):
                issues.append(_issue("invalid_image", "页面图片无法读取或文件已损坏。", slide_id))
                continue
            if width <= 0 or height <= 0 or abs((width / height) - ASPECT_RATIO) > ASPECT_TOLERANCE:
                issues.append(
                    _issue(
                        "invalid_aspect_ratio",
                        f"页面图片必须是 16:9，当前为 {width}×{height}。",
                        slide_id,
                    )
                )
                continue
            provenance = visual_provenance_status(root, slide_id)
            if not provenance.get("valid"):
                issues.append(
                    _issue(
                        "stale_or_unconfirmed_image",
                        "页面图片未确认、来源记录缺失或已因分镜变化而失效。",
                        slide_id,
                    )
                )
                continue
            slides.append(
                {
                    "slide_id": slide_id,
                    "image_path": str(image_path),
                    "width": width,
                    "height": height,
                }
            )
    return {
        "ready": not issues and bool(slide_ids),
        "export_type": "pptx",
        "format": "image_only",
        "slide_count": len(slide_ids),
        "ready_slide_count": len(slides),
        "slide_ids": slide_ids,
        "slides": slides,
        "issues": issues,
    }


def _validate_saved_presentation(path: Path, expected_slide_count: int) -> None:
    presentation = Presentation(str(path))
    if len(presentation.slides) != expected_slide_count:
        raise RuntimeError(
            f"PPTX 页数校验失败：预期 {expected_slide_count} 页，实际 {len(presentation.slides)} 页"
        )
    if presentation.slide_width != Inches(SLIDE_WIDTH_INCHES):
        raise RuntimeError("PPTX 页面宽度校验失败")
    if presentation.slide_height != Inches(SLIDE_HEIGHT_INCHES):
        raise RuntimeError("PPTX 页面高度校验失败")
    for index, slide in enumerate(presentation.slides, start=1):
        if not any(shape.shape_type == 13 for shape in slide.shapes):
            raise RuntimeError(f"PPTX 第 {index} 页缺少图片")


def build_image_only_pptx(
    run_dir: str | Path,
    filename: str,
    *,
    title: str = "",
    progress: Callable[[int, str], None] | None = None,
) -> dict[str, Any]:
    root = Path(run_dir).resolve()
    readiness = inspect_pptx_readiness(root)
    if not readiness["ready"]:
        raise PptxReadinessError(readiness)
    source_fingerprint = presentation_input_fingerprint(root)
    progress = progress or (lambda _value, _stage: None)
    target_dir = presentations_dir(root)
    target_dir.mkdir(parents=True, exist_ok=True)
    target = presentation_file(root, filename)
    temporary = Path(f"{target}.part")
    if temporary.exists():
        temporary.unlink()

    progress(20, "composing")
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
    presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    presentation.core_properties.title = str(title or "PPT Studio 图片演示文稿")
    presentation.core_properties.subject = "由已确认的 Slide 图片生成"
    presentation.core_properties.author = "PPT Visualization Studio"
    blank_layout = presentation.slide_layouts[6]
    for index, slide_data in enumerate(readiness["slides"], start=1):
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            slide_data["image_path"],
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
        # 将该页演讲稿写入演讲者备注（notes），便于放映时用演讲者视图逐页讲解。
        # 备注内容不遮挡图片：备注属于独立的 notes 区，仅演讲者可见。
        slide_id = slide_data["slide_id"]
        notes_text = _read_slide_notes(root, slide_id)
        if notes_text:
            notes = slide.notes_slide
            notes.notes_text_frame.text = notes_text
        progress(20 + round(60 * index / len(readiness["slides"])), "composing")

    try:
        presentation.save(str(temporary))
        progress(88, "verifying")
        _validate_saved_presentation(temporary, readiness["slide_count"])
        os.replace(temporary, target)
    finally:
        if temporary.exists():
            temporary.unlink()

    fingerprint_after_export = presentation_input_fingerprint(root)
    if fingerprint_after_export.get("digest") != source_fingerprint.get("digest"):
        if target.exists():
            target.unlink()
        raise RuntimeError("生成过程中分镜或页面图片发生变化，请重新生成 PPTX")
    exported_at = datetime.now().isoformat(timespec="seconds")
    metadata = {
        "schema_version": 1,
        "export_version": PPTX_EXPORT_VERSION,
        "exported_at": exported_at,
        "filename": filename,
        "slide_count": readiness["slide_count"],
        "slide_ids": readiness["slide_ids"],
        "slide_size": {
            "width_inches": SLIDE_WIDTH_INCHES,
            "height_inches": SLIDE_HEIGHT_INCHES,
            "aspect_ratio": "16:9",
        },
        "content_mode": "full_slide_bitmap",
        "source_fingerprint": source_fingerprint,
        "output_sha256": sha256_file(target),
    }
    try:
        write_json_atomic(presentation_sidecar(target), metadata)
    except Exception:
        if target.exists():
            target.unlink()
        raise
    progress(100, "completed")
    return {
        "path": str(target),
        "filename": filename,
        "size_bytes": target.stat().st_size,
        "mime_type": PPTX_MIME_TYPE,
        "fingerprint": source_fingerprint,
        "metadata": metadata,
    }
