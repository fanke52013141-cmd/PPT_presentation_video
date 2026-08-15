"""Reveal-style PPT export: decompose each slide into N+1 stacked pages by mask order.

For each source slide:
  page 0: base image (master with all dynamic mask regions erased to background)
  page K: base + first K dynamic crops overlaid in z_index order (PNG alpha keeps transparency)

This produces a PowerPoint that plays like an animation: elements appear one by one
as the viewer advances pages, matching the reveal experience users get from the video
pipeline but without needing Remotion.
"""

from __future__ import annotations

import json
import logging
import os
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.util import Inches, Emu

from artifact_fingerprint import presentation_input_fingerprint, sha256_file
from pipeline_lifecycle import write_json_atomic
from project_storage import (
    UnsafeProjectPath,
    presentation_file,
    presentation_sidecar,
    presentations_dir,
    safe_identifier,
    slide_dir,
    slide_file,
    visual_contract_path,
)

# 复用视频侧的 mask 光栅化 + 抠图工具
try:
    from scripts.background_color import masked_outer_white_cutout
    from scripts.build_reveal_scene import (
        DEFAULT_CANVAS,
        RevealBuildError,
        alpha_box,
        hex_to_rgb,
        manual_mask_alpha,
        read_json,
    )
except ModuleNotFoundError:
    from background_color import masked_outer_white_cutout
    from build_reveal_scene import (
        DEFAULT_CANVAS,
        RevealBuildError,
        alpha_box,
        hex_to_rgb,
        manual_mask_alpha,
        read_json,
    )


logger = logging.getLogger("PPTStudio.PptxRevealExport")

PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PPTX_REVEAL_EXPORT_VERSION = "reveal_pptx_v1"
SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5
DEFAULT_BACKGROUND_HEX = "#FEFDF9"


class PptxRevealExportError(RuntimeError):
    """Raised when reveal-style export cannot proceed (manifest missing, mask invalid, etc.).**

    status_code lets the route layer map to HTTP responses uniformly.
    """

    def __init__(self, detail: str, status_code: int = 400) -> None:
        super().__init__(detail)
        self.detail = detail
        self.status_code = status_code


def _utc_now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def _load_reveal_manifest(run_dir: Path) -> dict[str, Any]:
    """Read and validate reveal_manifest.json under run_dir."""
    manifest_path = run_dir / "reveal_manifest.json"
    if not manifest_path.exists():
        raise PptxRevealExportError(
            "缺少 reveal_manifest.json，请先在『Mask 标注』步骤完成至少一次自动或手动标注。",
            status_code=409,
        )
    try:
        manifest = read_json(manifest_path)
    except Exception as exc:
        raise PptxRevealExportError(
            f"reveal_manifest.json 读取失败：{exc}",
            status_code=500,
        ) from exc
    version = str(manifest.get("version") or "")
    if version != "reveal_v1":
        raise PptxRevealExportError(
            f"reveal_manifest.json 版本不支持：{version!r}，预期 'reveal_v1'",
            status_code=409,
        )
    return manifest


def _slide_manifest(manifest: dict[str, Any], slide_id: str) -> dict[str, Any] | None:
    for slide in manifest.get("slides", []) or []:
        if slide.get("slide_id") == slide_id:
            return slide
    return None


def _collect_dynamic_groups(slide_manifest: dict[str, Any]) -> list[dict[str, Any]]:
    """Return groups that have a non-empty manual_mask and are not flagged static, sorted by z_index asc."""
    groups = slide_manifest.get("groups") or []
    dynamic = []
    for group in groups:
        manual_mask = group.get("manual_mask")
        if not manual_mask:
            continue
        if group.get("is_static") or group.get("is_static_header"):
            continue
        rle = (manual_mask.get("rle") or {}) if isinstance(manual_mask, dict) else {}
        runs = rle.get("runs") if isinstance(rle, dict) else None
        strokes = manual_mask.get("strokes") if isinstance(manual_mask, dict) else None
        if not runs and not strokes:
            continue
        dynamic.append(group)
    dynamic.sort(key=lambda g: g.get("z_index", 0))
    return dynamic


def _build_base_image(
    master_rgb: Image.Image,
    dynamic_alphas: list[Image.Image],
    background_rgb: tuple[int, int, int],
) -> Image.Image:
    """Erase all dynamic mask regions from master, replacing them with the background color.

    Result: '非 mask 标注区域' 的原貌——未画 mask 的元素全部保留在底图里。
    """
    width, height = master_rgb.size
    if not dynamic_alphas:
        return master_rgb.copy()
    union_alpha = Image.new("L", (width, height), 0)
    for alpha in dynamic_alphas:
        union_alpha = ImageChops.lighter(union_alpha, alpha)
    background_layer = Image.new("RGB", (width, height), background_rgb)
    base = master_rgb.copy()
    base.paste(background_layer, (0, 0), mask=union_alpha)
    return base


def _pixels_to_inches(px: float, reference_px: int, reference_in: float) -> float:
    return px * reference_in / reference_px


def _read_slide_notes(run_dir: Path, slide_id: str) -> str:
    notes_path = slide_dir(run_dir, slide_id) / "narration.txt"
    if notes_path.exists():
        try:
            return notes_path.read_text(encoding="utf-8")
        except Exception:
            return ""
    return ""


def inspect_reveal_pptx_readiness(run_dir: str | Path) -> dict[str, Any]:
    """Check whether a project is ready for reveal-style PPT export."""
    root = Path(run_dir)
    vc_path = visual_contract_path(root)
    slide_ids: list[str] = []
    if vc_path.exists():
        try:
            vc = json.loads(vc_path.read_text(encoding="utf-8"))
            for slide in vc.get("slides", []) or []:
                sid = slide.get("slide_id")
                if sid:
                    slide_ids.append(str(sid))
        except Exception:
            pass

    issues: list[dict[str, str]] = []
    try:
        manifest = _load_reveal_manifest(root)
    except PptxRevealExportError as exc:
        issues.append({"code": "manifest_missing", "message": exc.detail})
        manifest = None

    canvas = (manifest or {}).get("canvas") or {}
    canvas_width = int(canvas.get("width") or DEFAULT_CANVAS["width"])
    canvas_height = int(canvas.get("height") or DEFAULT_CANVAS["height"])
    background_hex = str(canvas.get("background") or DEFAULT_BACKGROUND_HEX)

    slides_info: list[dict[str, Any]] = []
    total_pages = 0
    if manifest:
        for slide_id in slide_ids:
            slide_manifest = _slide_manifest(manifest, slide_id)
            if not slide_manifest:
                issues.append({
                    "code": "slide_not_in_manifest",
                    "message": f"{slide_id} 不在 reveal_manifest.json 中",
                    "slide_id": slide_id,
                })
                continue
            master_path = slide_dir(root, slide_id) / (slide_manifest.get("master") or "visual_draft.png")
            if not master_path.exists():
                issues.append({
                    "code": "master_missing",
                    "message": f"{slide_id} 缺少底图 {master_path.name}",
                    "slide_id": slide_id,
                })
                continue
            dynamic_groups = _collect_dynamic_groups(slide_manifest)
            slides_info.append({
                "slide_id": slide_id,
                "image_path": str(master_path),
                "width": canvas_width,
                "height": canvas_height,
                "dynamic_group_count": len(dynamic_groups),
                "page_count": len(dynamic_groups) + 1,
            })
            total_pages += len(dynamic_groups) + 1

    return {
        "ready": bool(slide_ids) and not issues and bool(slides_info),
        "export_type": "pptx",
        "format": "reveal",
        "slide_count": len(slide_ids),
        "ready_slide_count": len(slides_info),
        "total_page_count": total_pages,
        "slide_ids": slide_ids,
        "canvas": {"width": canvas_width, "height": canvas_height, "background": background_hex},
        "slides": slides_info,
        "issues": issues,
    }


def build_reveal_pptx(
    run_dir: str | Path,
    filename: str,
    *,
    title: str = "",
    progress: Callable[[int, str], None] | None = None,
) -> dict[str, Any]:
    """Build a reveal-style .pptx where each slide is expanded into N+1 stacked pages.

    Raises PptxRevealExportError when manifest/mask data is missing or invalid.
    """
    root = Path(run_dir)
    readiness = inspect_reveal_pptx_readiness(root)
    if not readiness["ready"]:
        messages = "; ".join(issue["message"] for issue in readiness["issues"]) or "未就绪"
        raise PptxRevealExportError(messages, status_code=409)

    manifest = _load_reveal_manifest(root)
    canvas = manifest.get("canvas") or {}
    canvas_width = int(canvas.get("width") or DEFAULT_CANVAS["width"])
    canvas_height = int(canvas.get("height") or DEFAULT_CANVAS["height"])
    background_hex = str(canvas.get("background") or DEFAULT_BACKGROUND_HEX)
    background_rgb = hex_to_rgb(background_hex)
    slide_ids: list[str] = readiness["slide_ids"]

    target_dir = presentations_dir(root)
    target_dir.mkdir(parents=True, exist_ok=True)
    target = presentation_file(root, filename)
    part_path = target.with_suffix(target.suffix + ".part")

    iso = _utc_now_iso()
    if progress:
        progress(5, "准备画布")

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
    presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    blank_layout = presentation.slide_layouts[6]

    inch_per_px_x = SLIDE_WIDTH_INCHES / canvas_width
    inch_per_px_y = SLIDE_HEIGHT_INCHES / canvas_height

    slide_step = max(1, len(slide_ids))
    progress_pct = 10
    notes_for_slides: dict[str, str] = {}

    for slide_index, slide_id in enumerate(slide_ids):
        slide_manifest = _slide_manifest(manifest, slide_id)
        if not slide_manifest:
            raise PptxRevealExportError(
                f"{slide_id} 不在 reveal_manifest.json 中",
                status_code=409,
            )

        master_path = slide_dir(root, slide_id) / (slide_manifest.get("master") or "visual_draft.png")
        if not master_path.exists():
            raise PptxRevealExportError(
                f"{slide_id} 缺少底图 {master_path.name}",
                status_code=409,
            )

        try:
            master_image = Image.open(master_path).convert("RGB")
        except Exception as exc:
            raise PptxRevealExportError(
                f"{slide_id} 底图读取失败：{exc}",
                status_code=500,
            ) from exc

        if master_image.size != (canvas_width, canvas_height):
            master_image = master_image.resize((canvas_width, canvas_height))

        dynamic_groups = _collect_dynamic_groups(slide_manifest)

        if progress:
            progress(
                min(90, progress_pct),
                f"{slide_id}：光栅化 {len(dynamic_groups)} 个 mask",
            )

        crops: list[tuple[Image.Image, dict[str, int]]] = []
        dynamic_alphas: list[Image.Image] = []
        for group in dynamic_groups:
            try:
                alpha = manual_mask_alpha(group.get("manual_mask"), canvas_width, canvas_height)
            except RevealBuildError as exc:
                raise PptxRevealExportError(
                    f"{slide_id} 的 mask {group.get('id', '?')} 解析失败：{exc}",
                    status_code=409,
                ) from exc
            if alpha is None:
                continue
            dynamic_alphas.append(alpha)
            try:
                layer, _, _ = masked_outer_white_cutout(master_image, alpha)
            except Exception as exc:
                raise PptxRevealExportError(
                    f"{slide_id} 的 mask {group.get('id', '?')} 抠图失败：{exc}",
                    status_code=500,
                ) from exc
            box = alpha_box(alpha, canvas_width, canvas_height)
            crops.append((layer, box))

        base_image = _build_base_image(master_image, dynamic_alphas, background_rgb)
        base_png_path = slide_dir(root, slide_id) / "pptx_reveal_base.png"
        try:
            base_image.save(base_png_path, format="PNG")
        except Exception as exc:
            raise PptxRevealExportError(
                f"{slide_id} 底图保存失败：{exc}",
                status_code=500,
            ) from exc

        crop_paths: list[Path] = []
        for idx, (layer, _box) in enumerate(crops):
            crop_path = slide_dir(root, slide_id) / f"pptx_reveal_crop_{idx:02d}.png"
            try:
                layer.save(crop_path, format="PNG")
            except Exception as exc:
                raise PptxRevealExportError(
                    f"{slide_id} 第 {idx + 1} 个 crop 保存失败：{exc}",
                    status_code=500,
                ) from exc
            crop_paths.append(crop_path)

        notes_text = _read_slide_notes(root, slide_id)
        if notes_text:
            notes_for_slides[slide_id] = notes_text

        page_count = len(crop_paths) + 1
        for k in range(page_count):
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(base_png_path),
                0,
                0,
                width=Inches(SLIDE_WIDTH_INCHES),
                height=Inches(SLIDE_HEIGHT_INCHES),
            )
            for idx in range(k):
                crop_path = crop_paths[idx]
                _layer, box = crops[idx]
                left_in = _pixels_to_inches(box["x"], canvas_width, SLIDE_WIDTH_INCHES)
                top_in = _pixels_to_inches(box["y"], canvas_height, SLIDE_HEIGHT_INCHES)
                width_in = _pixels_to_inches(box["w"], canvas_width, SLIDE_WIDTH_INCHES)
                height_in = _pixels_to_inches(box["h"], canvas_height, SLIDE_HEIGHT_INCHES)
                slide.shapes.add_picture(
                    str(crop_path),
                    Inches(left_in),
                    Inches(top_in),
                    width=Inches(width_in),
                    height=Inches(height_in),
                )
            if k == page_count - 1 and notes_text:
                slide.notes_slide.notes_text_frame.text = notes_text

        if progress:
            progress_pct = 10 + int((slide_index + 1) / slide_step * 80)
            progress(min(90, progress_pct), f"{slide_id}：已生成 {page_count} 页")

        try:
            base_png_path.unlink(missing_ok=True)
        except Exception:
            pass
        for crop_path in crop_paths:
            try:
                crop_path.unlink(missing_ok=True)
            except Exception:
                pass

    if progress:
        progress(95, "保存 PPTX")

    presentation.save(str(part_path))

    try:
        from pptx import Presentation as _P
        _verify = _P(str(part_path))
        verify_slide_count = len(_verify.slides)
    except Exception as exc:
        try:
            part_path.unlink(missing_ok=True)
        except Exception:
            pass
        raise PptxRevealExportError(
            f"PPTX 校验失败：{exc}",
            status_code=500,
        ) from exc

    expected_pages = sum(slide["page_count"] for slide in readiness["slides"])
    if verify_slide_count != expected_pages:
        try:
            part_path.unlink(missing_ok=True)
        except Exception:
            pass
        raise PptxRevealExportError(
            f"PPTX 页数校验失败：期望 {expected_pages}，实际 {verify_slide_count}",
            status_code=500,
        )

    os.replace(part_path, target)

    if progress:
        progress(98, "写元数据")

    size_bytes = target.stat().st_size
    output_sha = sha256_file(target)
    source_fp = presentation_input_fingerprint(root)
    metadata = {
        "schema_version": 1,
        "export_version": PPTX_REVEAL_EXPORT_VERSION,
        "exported_at": iso,
        "filename": filename,
        "slide_count": len(slide_ids),
        "slide_ids": slide_ids,
        "total_page_count": verify_slide_count,
        "slide_size": {
            "width_inches": SLIDE_WIDTH_INCHES,
            "height_inches": SLIDE_HEIGHT_INCHES,
            "aspect_ratio": "16:9",
        },
        "canvas": {
            "width": canvas_width,
            "height": canvas_height,
            "background": background_hex,
        },
        "content_mode": "reveal_layers",
        "reveal_pipeline_version": "exact_rle_mask_with_manual_corrections_v5",
        "source_fingerprint": source_fp,
        "output_sha256": output_sha,
        "notes_included_slide_ids": list(notes_for_slides.keys()),
    }
    try:
        write_json_atomic(presentation_sidecar(target), metadata)
    except Exception as exc:
        try:
            target.unlink(missing_ok=True)
        except Exception:
            pass
        raise PptxRevealExportError(
            f"元数据写入失败：{exc}",
            status_code=500,
        ) from exc

    if progress:
        progress(100, "完成")

    return {
        "path": str(target),
        "filename": filename,
        "size_bytes": size_bytes,
        "mime_type": PPTX_MIME_TYPE,
        "fingerprint": source_fp,
        "metadata": metadata,
    }
